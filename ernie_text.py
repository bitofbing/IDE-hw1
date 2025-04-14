# -*- coding: utf-8 -*-
import re
import jieba.posseg as pseg
from collections import defaultdict
from pptx import Presentation
from pathlib import Path
import paddle
import numpy as np
from sklearn.cluster import KMeans
from paddlenlp.transformers import ErnieForSequenceClassification, ErnieTokenizer

# 加载ERNIE模型和tokenizer
model_name = "ernie-3.0-medium-zh"
tokenizer = ErnieTokenizer.from_pretrained(model_name)

# 领域关键词增强版
DOMAIN_KEYWORDS = {
    "技术概念": ["优化", "查询", "数据库", "分布式", "架构", "向量", "基数", "索引"],
    "方法/模型": ["直方图", "采样", "回归", "神经网络", "贝叶斯", "卷积", "图谱"],
    "工具/系统": ["Spark", "Flink", "PostgreSQL", "CRF", "BMES"],
    "组织/机构": ["大学", "实验室", "研究院", "出版社"],
    "性能指标": ["吞吐量", "延迟", "准确率", "NDV"]
}

class EntityRelationExtractor:
    def __init__(self):
        self.term_db = defaultdict(set)
        self._init_term_db()  # 初始化术语库
        self.model = ErnieForSequenceClassification.from_pretrained(
            model_name,
            num_classes=len(DOMAIN_KEYWORDS)+1)
        self.model.eval()

        self.relation_types = {
            "作用": ["用于", "应用于", "用来", "适用于"],
            "依赖": ["基于", "依赖于", "建立在"],
            "比较": ["对比", "相比", "相较于"],
            "影响": ["提高", "降低", "增加", "减少"],
            "组成": ["包含", "组成", "构成"],
            "同类": ["和", "与", "以及"]
        }

    def _init_term_db(self):
        """初始化基础术语库"""
        self.term_db.update({
            "技术概念": {"查询优化", "基数估计", "分布式计算"},
            "方法/模型": {"频率直方图", "混合直方图", "图神经网络"},
            "工具/系统": set(),
            "组织/机构": {"清华大学出版社"},
            "性能指标": set()
        })

    def update_term_db(self, text: str):
        """动态更新术语库"""
        words = [word for word, pos in pseg.cut(text) if pos.startswith('n')]

        candidates = set()
        for i in range(len(words) - 1):
            term = words[i] + words[i + 1]
            if self._is_valid_term(term):
                candidates.add(term)
            if i < len(words) - 2 and len(words[i + 1]) < 4:
                term3 = words[i] + words[i + 1] + words[i + 2]
                if self._is_valid_term(term3):
                    candidates.add(term3)

        for term in candidates:
            added = False
            for cat, keywords in DOMAIN_KEYWORDS.items():
                if any(kw in term for kw in keywords):
                    self.term_db[cat].add(term)
                    added = True
                    break
            if not added and any(t in term for t in ["网络", "算法"]):
                self.term_db["方法/模型"].add(term)

    def _is_valid_term(self, term: str) -> bool:
        """验证术语是否有效"""
        if len(term) < 2 or len(term) > 12:
            return False

        words = pseg.cut(term)
        pos_tags = [pos for word, pos in words]
        if not any(tag.startswith('n') for tag in pos_tags):
            return False

        for cat, keywords in DOMAIN_KEYWORDS.items():
            if any(kw in term for kw in keywords):
                return True
        return False

    def extract_entities(self, text: str) -> dict:
        """实体识别"""
        entities = defaultdict(list)

        # 1. 使用术语库匹配
        sorted_terms = sorted(
            [(cat, term) for cat, terms in self.term_db.items() for term in terms],
            key=lambda x: len(x[1]),
            reverse=True)

        matched_positions = set()
        for cat, term in sorted_terms:
            if term in text:
                start = 0
                while True:
                    idx = text.find(term, start)
                    if idx == -1:
                        break
                    if not any(idx <= p < idx + len(term) for p in matched_positions):
                        entities[cat].append((term, idx, idx + len(term)))
                        matched_positions.update(range(idx, idx + len(term)))
                    start = idx + len(term)

        # 2. 使用规则匹配组织/机构
        org_pattern = r"([\u4e00-\u9fa5]{2,10}(公司|集团|大学|学院|研究院|实验室|出版社))"
        for match in re.finditer(org_pattern, text):
            org_name = match.group(1)
            if any(kw in org_name for kw in DOMAIN_KEYWORDS["组织/机构"]):
                entities["组织/机构"].append((org_name, match.start(), match.end()))

        return dict(entities)

    def extract_relations(self, text: str, entities: dict) -> list:
        """关系抽取"""
        relations = []
        current_ents = {e[0] for ents in entities.values() for e in ents}

        patterns = [
            (r"([\w\u4e00-\u9fa5]{2,10})(主要用于|专门用于|适用于|用于)([\w\u4e00-\u9fa5]{2,10})", "作用"),
            (r"([\w\u4e00-\u9fa5]{2,10})(基于|依赖于|建立在)([\w\u4e00-\u9fa5]{2,10})(技术|方法|模型)?", "依赖"),
            (r"([\w\u4e00-\u9fa5]{2,10})(与|和)([\w\u4e00-\u9fa5]{2,10})(的)?(结合|对比|比较)", "关联"),
            (r"([\w\u4e00-\u9fa5]{2,10})(显著|明显)?(提高|降低|减少|增加)([\w\u4e00-\u9fa5]{2,10})", "影响"),
            (r"([\w\u4e00-\u9fa5]{2,10})是([\w\u4e00-\u9fa5]{2,10})(的)?(组成部分|关键要素)", "组成"),
            (r"([\w\u4e00-\u9fa5]{2,10})(在)([\w\u4e00-\u9fa5]{2,10})(中)?(扮演角色|起到作用|关键作用)?", "角色"),
            (r"([\w\u4e00-\u9fa5]{2,10})(对|作用于)([\w\u4e00-\u9fa5]{2,10})", "作用"),
            (r"([\w\u4e00-\u9fa5]{2,10})(与|和)([\w\u4e00-\u9fa5]{2,10})(有)?(关系)", "关系")
        ]

        for pattern, rel_type in patterns:
            for match in re.finditer(pattern, text):
                print(f"匹配到关系: {match.group(1)} → {rel_type} → {match.group(3)}")  # 输出匹配信息
                subj, obj = match.group(1), match.group(3)
                if subj in current_ents and obj in current_ents:
                    relations.append((subj, rel_type, obj))

        return relations

def extract_ppt_text(file_path: str) -> list:
    """提取PPT文本"""
    prs = Presentation(file_path)
    return [
        (i + 1, "\n".join(shape.text.strip()
                          for shape in slide.shapes
                          if hasattr(shape, "text") and shape.text.strip()))
        for i, slide in enumerate(prs.slides)
    ]

def process_all_ppts(folder_path: str = "/home/aistudio/data") -> dict:
    """处理流程"""
    extractor = EntityRelationExtractor()
    results = {}

    for file_path in Path(folder_path).glob("*.pptx"):
        file_results = []
        print(f"\n处理文件: {file_path.name}")

        for page_num, text in extract_ppt_text(file_path):
            extractor.update_term_db(text)  # 更新术语库
            entities = extractor.extract_entities(text)
            relations = extractor.extract_relations(text, entities)

            file_results.append({
                "page": page_num,
                "text": text[:100] + "..." if len(text) > 100 else text,
                "entities": {k: [e[0] for e in v] for k, v in entities.items()},
                "relations": relations
            })

        results[file_path.name] = {
            "slides": file_results,
            "final_terms": {k: list(v) for k, v in extractor.term_db.items()}
        }

    return results

def print_results(results: dict):
    """结果显示"""
    for file_name, data in results.items():
        print(f"\n\n=== 文件 {file_name} ===")
        print("最终术语库:")
        for cat, terms in data["final_terms"].items():
            if terms:
                print(f"  {cat}:")
                for i in range(0, len(terms), 5):
                    print("   ", ", ".join(terms[i:i + 5]))

        for slide in data["slides"]:
            if not slide["entities"] and not slide["relations"]:
                continue

            print(f"\nSlide {slide['page']}:")
            if slide["entities"]:
                print("实体:")
                for cat, ents in slide["entities"].items():
                    if ents:
                        print(f"  {cat}: {', '.join(ents)}")
            if slide["relations"]:
                print("关系:")
                for rel in slide["relations"]:
                    print(f"  {rel[0]} → {rel[1]} → {rel[2]}")

if __name__ == "__main__":
    paddle.set_device('gpu')
    print("ERNIE模型加载完成，开始处理PPT文件...")
    analysis_results = process_all_ppts()
    # print_results(analysis_results)